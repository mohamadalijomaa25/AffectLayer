import os

def generate_docx():
    try:
        from docx import Document
    except ImportError:
        print("Please install python-docx: pip install python-docx")
        return

    doc = Document()
    doc.add_heading('Project Report: AffectLayer', 0)
    
    doc.add_heading('1. Introduction', level=1)
    doc.add_heading('Problem Statement', level=2)
    doc.add_paragraph('Standard sentiment analysis algorithms evaluate plain text at face value. When a user utilizes masked language—such as passive-aggression or minimization—to hide their emotional distress, traditional models mislabel the context as "positive" or "neutral". There is a gap in computational linguistics for models that identify explicitly polite phrasing acting as a psychological coping mechanism.')
    doc.add_heading('Motivation', level=2)
    doc.add_paragraph('By identifying latent emotional states, we can improve mental health AI monitors, social media sentiment trackers, and create more empathetic human-computer interaction designs.')
    doc.add_heading('Project Objective', level=2)
    doc.add_paragraph('Our objective was to build a Tri-Engine NLP application demonstrating the architectural differences between Classical NLP rules, Machine Learning embeddings, and Generative LLM logic when analyzing complex human subtext.')

    doc.add_heading('2. Methodology', level=1)
    doc.add_heading('Dataset & Preprocessing', level=2)
    doc.add_paragraph('Because AffectLayer is an architectural evaluation, we utilized a curated "Validation Dataset" of 15 conversational phrases deeply embedded with emotional masking. Preprocessing is handled dynamically on the backend, converting all raw inputs into acceptable JSON structures for the LLMs.')
    doc.add_heading('Models / Tools Used', level=2)
    doc.add_paragraph('1. VADER (Classical NLP): A deterministic, rule-based lexicon.')
    doc.add_paragraph('2. DistilBERT (Machine Learning): A neural network fine-tuned on SST-2, accessed via the Hugging Face Serverless API.')
    doc.add_paragraph('3. Gemini 2.5 (Generative LLM): Google\'s large language model accessed via zero-shot prompt engineering.')
    doc.add_paragraph('The frontend was built with React, and the backend relay was built with Python (FastAPI).')
    
    doc.add_heading('3. Results', level=1)
    doc.add_heading('Analysis & Comparison of Methods', level=2)
    doc.add_paragraph('Our experiments proved that the Classical Engine fails immediately when attempting to detect sarcasm, often generating false positives. The ML DistilBERT engine understands contextual negativity better but fails to diagnose specific psychological mechanisms. The Generative LLM (Gemini) proved vastly superior, successfully identifying mechanisms like "Evasion" and "Passive Aggression", whilst generating accurate clinical rationale in both English and Arabic.')
    doc.add_heading('Limitations', level=2)
    doc.add_paragraph('LLM and ML engines suffer from slight API latency compared to the instant calculation of the Classical engine, making real-time analysis heavier on network load.')

    doc.add_heading('4. Conclusion', level=1)
    doc.add_paragraph('AffectLayer successfully bridges the gap between literal sentiment analysis and complex psychological extraction, proving that next-generation generative AI is the only system currently capable of robust "masked emotion" detection.')

    output_path = os.path.join(os.getcwd(), 'AffectLayer_Final_Report.docx')
    doc.save(output_path)
    print(f"Generated {output_path}")

def generate_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("Please install python-pptx: pip install python-pptx")
        return

    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AffectLayer: Hidden Emotion Detection"
    slide.placeholders[1].text = "Unmasking Latent Distress Using Tri-Engine NLP Architecture"

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Traditional NLP takes text too literally."
    p = tf.add_paragraph()
    p.text = "If someone says 'Oh perfect, exactly what I needed' after getting a flat tire, basic AI thinks they are happy."
    p = tf.add_paragraph()
    p.text = "There is a gap in detecting 'Emotionally Masked' language."
    
    # Slide 3: The Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Tri-Engine Methodology"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "We compare three distinct AI generations:"
    p = tf.add_paragraph()
    p.text = "1. Classical NLP (VADER): Fast but literal."
    p = tf.add_paragraph()
    p.text = "2. Machine Learning (DistilBERT via Hugging Face): Understands context but misses psychology."
    p = tf.add_paragraph()
    p.text = "3. Generative AI (Gemini 2.5): Deep psychological extraction with Dual-Language Rationales."

    # Slide 4: Demo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Live Demonstration"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Switching to the AffectLayer platform preview..."
    p = tf.add_paragraph()
    p.text = "We will test the sentence: 'Haha it's okay, I'm used to being ignored anyway.'"
    
    # Slide 5: Results & Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Results and Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Results Summary:"
    p = tf.add_paragraph()
    p.text = "Generative AI is currently the only engine capable of successfully classifying masking mechanisms (e.g. Sarcasm, Passive-Aggression)."
    p = tf.add_paragraph()
    p.text = "By separating our UI in React and our backend in Python FastAPI, we created a lightning fast, highly optimized web tool."

    output_path = os.path.join(os.getcwd(), 'AffectLayer_Presentation.pptx')
    prs.save(output_path)
    print(f"Generated {output_path}")

if __name__ == "__main__":
    print("Generating official documents...")
    generate_docx()
    generate_pptx()
    print("Done!")
